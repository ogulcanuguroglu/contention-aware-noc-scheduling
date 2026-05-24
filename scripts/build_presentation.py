# -*- coding: utf-8 -*-
"""Build CMP720 final presentation as an editable .pptx file."""

import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
PROJECT_ROOT = Path(__file__).resolve().parent.parent
CONTENT_MD   = PROJECT_ROOT / "cmp720_presentation_en.md"
OUTPUT_DIR   = PROJECT_ROOT / "presentation"
OUTPUT_PATH  = OUTPUT_DIR / "cmp720_final.pptx"
FIGURES_BASE = PROJECT_ROOT / "results" / "figures"

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------------------
# Colours & sizes
# ---------------------------------------------------------------------------
NAVY      = RGBColor(0x0B, 0x2E, 0x4F)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
BURGUNDY  = RGBColor(0x8B, 0x1A, 0x2B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY_BG = RGBColor(0xF2, 0xF2, 0xF2)
GREEN_TINT    = RGBColor(0xD6, 0xF0, 0xD6)
RED_TINT      = RGBColor(0xF7, 0xD6, 0xD6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_BAR_H  = Inches(0.55)
FOOTER_Y     = SLIDE_H - Inches(0.32)
CONTENT_TOP  = TITLE_BAR_H + Inches(0.12)
CONTENT_H    = SLIDE_H - CONTENT_TOP - Inches(0.38)

FONT_TITLE   = "Arial"
FONT_BODY    = "Arial"
FONT_MONO    = "Courier New"

# ---------------------------------------------------------------------------
# Figure catalogue  (slide_number → path or list of paths)
# ---------------------------------------------------------------------------
P21 = FIGURES_BASE / "phase21_interpretive"
P20 = FIGURES_BASE / "phase20_presentation"
P18 = FIGURES_BASE / "phase18"
CE  = FIGURES_BASE / "contention_explainer"

FIGURE_MAP = {
    1:  CE  / "figA_link_contention_mechanism.png",
    2:  P18 / "fig0a_noc_topology_generic.png",
    # 3: table slide — no figure
    4:  P21 / "fig3_native_vs_replay.png",
    5:  P21 / "fig2_scheduler_concept.png",
    6:  CE  / "figB_earliest_route_slot.png",
    # 7: text only
    8:  P21 / "fig1_dag_family_topologies.png",
    9:  [P21 / "fig5_ccr_sweep_replayed_speedup.png",
         P21 / "fig6_replay_overhead_ratio.png"],
    10: P21 / "fig4c_fork_join_gantt.png",
    11: P21 / "fig7_task_instance_ratio.png",
    # 12: conclusions — no figure
}

FIGURE_DESC = {
    1:  "Link contention mechanism diagram",
    2:  "4×4 NoC topology with XY routing",
    4:  "Native vs. Replayed makespan comparison",
    5:  "CA-D scheduler concept pipeline",
    6:  "Earliest route slot / ΔEFT comparison",
    8:  "4 DAG family topologies",
    9:  "CCR sweep replayed speedup + Replay overhead ratio",
    10: "Fork-Join HEFT vs CA-D Gantt chart",
    11: "Task instance ratio (duplication overhead)",
}

# ---------------------------------------------------------------------------
# Helpers: XML / formatting
# ---------------------------------------------------------------------------

def _set_cell_bg(cell, rgb: RGBColor):
    """Set solid background colour on a table cell via lxml."""
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def _para_space(para, before_pt=2, after_pt=2):
    pPr = para._p.get_or_add_pPr()
    pPr.set("spcBef", str(before_pt * 100))
    pPr.set("spcAft", str(after_pt * 100))


def _add_text_box(slide, left, top, width, height,
                  text="", font_name=FONT_BODY, font_size=18,
                  bold=False, colour=DARK_GREY, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = text
    r.font.name      = font_name
    r.font.size      = Pt(font_size)
    r.font.bold      = bold
    r.font.color.rgb = colour
    return txBox, tf


def _add_burgundy_bar(slide, page_num: int, title_text: str):
    """Thin burgundy title bar across the top."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, TITLE_BAR_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BURGUNDY
    bar.line.fill.background()

    # Title text (left-aligned)
    tf = bar.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title_text
    r.font.name      = FONT_TITLE
    r.font.size      = Pt(22)
    r.font.bold      = True
    r.font.color.rgb = WHITE

    # Page number (right side)
    pg_box, pg_tf = _add_text_box(
        slide,
        SLIDE_W - Inches(0.7), Inches(0.05),
        Inches(0.6), TITLE_BAR_H - Inches(0.1),
        text=str(page_num),
        font_size=12, bold=True, colour=WHITE
    )
    pg_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _add_footer(slide):
    _add_text_box(
        slide,
        Inches(0.15), FOOTER_Y,
        Inches(9), Inches(0.3),
        text="CMP720 — Contention-Aware DAG Scheduling on NoC",
        font_size=9, colour=RGBColor(0x88, 0x88, 0x88)
    )


def _set_notes(slide, notes_text: str):
    if not notes_text.strip():
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text.strip()


# ---------------------------------------------------------------------------
# Markdown parser helpers
# ---------------------------------------------------------------------------

def _parse_inline(text: str):
    """Return list of (text_str, bold, italic, mono) tuples from inline md."""
    parts = []
    pattern = re.compile(
        r'\*\*(.+?)\*\*'     # bold
        r'|\*(.+?)\*'         # italic
        r'|`(.+?)`'           # mono
    )
    last = 0
    for m in pattern.finditer(text):
        if m.start() > last:
            parts.append((text[last:m.start()], False, False, False))
        if m.group(1) is not None:
            parts.append((m.group(1), True, False, False))
        elif m.group(2) is not None:
            parts.append((m.group(2), False, True, False))
        elif m.group(3) is not None:
            parts.append((m.group(3), False, False, True))
        last = m.end()
    if last < len(text):
        parts.append((text[last:], False, False, False))
    return parts


def _add_run_inline(para, text_str, bold=False, italic=False, mono=False,
                    base_size=18, base_colour=DARK_GREY, first_run=False):
    """Add a styled run to a paragraph."""
    spans = _parse_inline(text_str)
    for span_text, span_bold, span_italic, span_mono in spans:
        r = para.add_run()
        r.text            = span_text
        r.font.size       = Pt(base_size)
        r.font.bold       = bold or span_bold
        r.font.italic     = italic or span_italic
        r.font.color.rgb  = base_colour
        r.font.name       = FONT_MONO if (mono or span_mono) else FONT_BODY


# ---------------------------------------------------------------------------
# Content bullet renderer
# ---------------------------------------------------------------------------

def _add_bullets_to_frame(tf, bullets, start_para_idx=0):
    """
    bullets: list of (level, text_str)
    level 0 = top-level bullet, level 1 = sub-bullet
    """
    for i, (level, text_str) in enumerate(bullets):
        if i == 0 and start_para_idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        pPr = para._p.get_or_add_pPr()
        pPr.set("lvl", str(level))

        indent_left = Inches(0.25 + level * 0.25)
        indent_hang = Inches(0.2)

        pPr.set("indent",    str(-int(indent_hang)))
        pPr.set("marL",      str(int(indent_left + indent_hang)))

        # Bullet marker via buChar
        buChar = pPr.find(qn("a:buChar"))
        if buChar is None:
            buNone = pPr.find(qn("a:buNone"))
            if buNone is not None:
                pPr.remove(buNone)
            buClr = etree.SubElement(pPr, qn("a:buClr"))
            srgb  = etree.SubElement(buClr, qn("a:srgbClr"))
            srgb.set("val", "8B1A2B")  # BURGUNDY
            buSz  = etree.SubElement(pPr, qn("a:buSzPct"))
            buSz.set("val", "100000")
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", "▸" if level == 0 else "–")

        _para_space(para, before_pt=3 if level == 0 else 1, after_pt=1)
        size = 17 if level == 0 else 15
        _add_run_inline(para, text_str, base_size=size)


# ---------------------------------------------------------------------------
# Figure insertion
# ---------------------------------------------------------------------------

def _placeholder_box(slide, left, top, width, height, desc):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    box.line.color.rgb      = BURGUNDY
    box.line.width          = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text            = f"[Figure pending: {desc}]"
    r.font.size       = Pt(13)
    r.font.italic     = True
    r.font.color.rgb  = RGBColor(0x77, 0x77, 0x77)
    r.font.name       = FONT_BODY
    return False   # not found


def _insert_figure(slide, path, left, top, width, height, desc):
    if path and Path(path).exists():
        try:
            slide.shapes.add_picture(str(path), left, top, width, height)
            return True
        except Exception:
            pass
    return _placeholder_box(slide, left, top, width, height, desc) or False


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_cover_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Navy background
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Burgundy accent strip (left side)
    accent = slide.shapes.add_shape(1, 0, 0, Inches(0.25), SLIDE_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = BURGUNDY
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.8), Inches(12), Inches(1.6)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text            = "Contention-Aware Task Duplication\nfor DAG Scheduling on NoC"
    r.font.name       = FONT_TITLE
    r.font.size       = Pt(36)
    r.font.bold       = True
    r.font.color.rgb  = WHITE

    # Subtitle line
    sub_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(3.7), Inches(10), Inches(0.5)
    )
    tf2 = sub_box.text_frame
    p2  = tf2.paragraphs[0]
    r2  = p2.add_run()
    r2.text            = "CMP720 Master Project  ·  Oğulcan Uğuroğlu"
    r2.font.name       = FONT_BODY
    r2.font.size       = Pt(20)
    r2.font.color.rgb  = RGBColor(0xCC, 0xCC, 0xCC)

    # Institution / date
    info_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(4.4), Inches(10), Inches(0.8)
    )
    tf3 = info_box.text_frame
    tf3.word_wrap = True
    p3  = tf3.paragraphs[0]
    r3  = p3.add_run()
    r3.text            = "Hacettepe University — Department of Computer Engineering\nMay 2026"
    r3.font.name       = FONT_BODY
    r3.font.size       = Pt(16)
    r3.font.color.rgb  = RGBColor(0xAA, 0xAA, 0xAA)


def build_content_slide(prs, page_num, title, bullets, notes="",
                        fig_path=None, fig_desc="", dual_fig=False,
                        fig_path2=None, right_ratio=0.45):
    """
    Generic content slide with optional right-panel figure(s).
    dual_fig=True: place two figures stacked vertically on the right.
    right_ratio: fraction of slide width allocated to figure panel.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_burgundy_bar(slide, page_num, title)
    _add_footer(slide)

    has_fig   = fig_path is not None
    left_frac = 1.0 - right_ratio if has_fig else 1.0

    text_left  = Inches(0.25)
    text_top   = CONTENT_TOP + Inches(0.05)
    text_w     = (SLIDE_W - Inches(0.35)) * left_frac
    text_h     = CONTENT_H - Inches(0.05)

    txBox = slide.shapes.add_textbox(text_left, text_top, text_w, text_h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    _add_bullets_to_frame(tf, bullets)

    # Figure panel
    report_found = None
    if has_fig:
        fig_left = text_left + text_w + Inches(0.1)
        fig_top  = CONTENT_TOP + Inches(0.05)
        fig_w    = SLIDE_W - fig_left - Inches(0.15)
        fig_h    = CONTENT_H - Inches(0.05)

        if dual_fig and fig_path2:
            half_h = (fig_h - Inches(0.1)) / 2
            r1 = _insert_figure(slide, fig_path,  fig_left, fig_top,
                                 fig_w, half_h, fig_desc)
            r2 = _insert_figure(slide, fig_path2, fig_left,
                                 fig_top + half_h + Inches(0.1),
                                 fig_w, half_h, fig_desc + " (2)")
            report_found = r1 and r2
        else:
            report_found = _insert_figure(slide, fig_path,
                                          fig_left, fig_top, fig_w, fig_h,
                                          fig_desc)

    _set_notes(slide, notes)
    return report_found


def build_table_slide(prs, page_num, title, intro_bullets, table_data, notes=""):
    """
    Slide with a short intro text block and a native pptx table.
    table_data: list of rows; first row is header.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_burgundy_bar(slide, page_num, title)
    _add_footer(slide)

    # Intro bullets (compact, top portion)
    intro_h = Inches(1.1)
    txBox = slide.shapes.add_textbox(
        Inches(0.25), CONTENT_TOP + Inches(0.05),
        SLIDE_W - Inches(0.5), intro_h
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    _add_bullets_to_frame(tf, intro_bullets)

    # Table
    rows = len(table_data)
    cols = len(table_data[0])
    tbl_top  = CONTENT_TOP + intro_h + Inches(0.15)
    tbl_left = Inches(0.3)
    tbl_w    = SLIDE_W - Inches(0.6)
    tbl_h    = CONTENT_H - intro_h - Inches(0.25)

    table = slide.shapes.add_table(rows, cols, tbl_left, tbl_top,
                                   tbl_w, tbl_h).table

    col_widths = [Inches(2.4), Inches(1.8), Inches(2.2), Inches(5.5)]
    for ci, w in enumerate(col_widths[:cols]):
        table.columns[ci].width = w

    for ri, row_data in enumerate(table_data):
        is_header  = (ri == 0)
        is_even    = (ri % 2 == 0) and not is_header
        row_bg     = BURGUNDY if is_header else (LIGHT_GREY_BG if is_even else WHITE)
        text_col   = WHITE    if is_header else DARK_GREY

        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            _set_cell_bg(cell, row_bg)

            tf  = cell.text_frame
            tf.word_wrap = True
            p   = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            is_yes = (cell_text.strip().lower() == "yes")
            is_no  = (cell_text.strip().lower() == "no")

            if is_yes and not is_header:
                _set_cell_bg(cell, GREEN_TINT)
                text_col2 = RGBColor(0x1A, 0x6B, 0x1A)
            elif is_no and not is_header:
                _set_cell_bg(cell, RED_TINT)
                text_col2 = RGBColor(0x8B, 0x1A, 0x1A)
            else:
                text_col2 = text_col

            r = p.add_run()
            # Strip leading ** for proposed row
            clean = cell_text.strip().lstrip("*").rstrip("*")
            r.text            = clean
            r.font.name       = FONT_BODY
            r.font.size       = Pt(14)
            r.font.bold       = is_header or cell_text.startswith("**")
            r.font.color.rgb  = text_col2

    _set_notes(slide, notes)


# ---------------------------------------------------------------------------
# Slide data  (all 12 content slides + cover)
# ---------------------------------------------------------------------------

SLIDES = [
    # ── Slide 1 ──────────────────────────────────────────────────────────
    dict(
        slide_num=1,
        title="Problem Definition — Network Contention in DAG Scheduling",
        bullets=[
            (0, "**The Illusion of Ideal Networks:** Classical DAG schedulers (e.g., HEFT) calculate inter-processor communication costs analytically, ignoring physical network traffic and link sharing."),
            (0, "**The Reality of Contention:** In 2D-Mesh NoC architectures, parallel tasks routing data across shared links create severe queuing delays — native makespan predictions are highly optimistic."),
            (0, "**The Core Research Question:** \"Should we wait for data through a congested multi-hop path, or is it more efficient to duplicate the predecessor task locally?\""),
            (0, "**Our Core Focus:** CA-D (Contention-Aware Duplication) — an offline scheduling heuristic that models physical NoC link reservations to make traffic-aware duplication decisions."),
        ],
        notes="In this project we solve the problem of statically assigning workloads on multi-core NoC systems. Algorithms like HEFT assume the network is always idle. We aim to close this optimism gap by modeling real-world network contention.",
        fig=FIGURE_MAP[1], fig_desc=FIGURE_DESC[1], right_ratio=0.40,
    ),
    # ── Slide 2 ──────────────────────────────────────────────────────────
    dict(
        slide_num=2,
        title="Target Framework & System Models",
        bullets=[
            (0, "**Application Model (DAG):** Directed Acyclic Graphs — vertices represent computation costs; directed edges define data dependency volumes."),
            (0, "**CCR Spectrum:** Evaluated across 0.1 (compute-bound) to 10.0 (communication-bound)."),
            (0, "**Platform Model:** 4×4 homogeneous 2D Mesh Network-on-Chip — 16 processing tiles."),
            (0, "**Deterministic XY Routing:** Packets move strictly along X-axis first, then Y-axis."),
            (0, "**Link-Level Slot Reservation:** Analytical delays replaced by directional, atomic interval reservations on physical NoC link channels."),
        ],
        notes="Instead of multiplying data transfers by simple coefficients, we process them as timed interval reservations on directional links along the XY routing path.",
        fig=FIGURE_MAP[2], fig_desc=FIGURE_DESC[2], right_ratio=0.38,
    ),
    # ── Slide 3 ── TABLE SLIDE ────────────────────────────────────────────
    dict(
        slide_num=3,
        title="Implementation Status — The 4-Scheduler Evaluation",
        type="table",
        intro_bullets=[
            (0, "**Custom Python Simulation Framework:** Lightweight, reservation-based hardware execution simulator optimised for heavy state-cloning."),
            (0, "A 2×2 design matrix isolates the independent contributions of *task duplication* and *contention-awareness*."),
        ],
        table_data=[
            ["Scheduler", "Task Duplication", "NoC Contention Model", "Role"],
            ["HEFT Baseline",     "No",  "No",  "Contention-blind, duplication-free list scheduler"],
            ["CA-LS",             "No",  "Yes", "No duplication; strict physical NoC reservations"],
            ["CD-LS",             "Yes", "No",  "Parent task replication; contention-blind network"],
            ["**CA-D (Proposed)", "Yes", "Yes", "Co-optimises link contention + recursive duplication"],
        ],
        notes="We developed 4 schedulers on the same infrastructure so we can isolate the effect of both duplication and contention-awareness independently.",
    ),
    # ── Slide 4 ──────────────────────────────────────────────────────────
    dict(
        slide_num=4,
        title="Architectural Deviations & Design Decisions",
        bullets=[
            (0, "**Deviation 1: Custom Interval Simulator over Cycle-Accurate Tools (BookSim)**"),
            (1, "Required for sub-microsecond state cloning during scheduler candidate exploration."),
            (0, "**Decision 2: The \"Fair Replay\" Module**"),
            (1, "Contention-blind schedulers produce falsely optimistic end times."),
            (1, "A validation engine re-enforces all placements on the same physical NoC model to extract true makespans."),
            (0, "**Decision 3: Structural Shallow-Copy Optimisation**"),
            (1, "Replaced `deepcopy` with shared immutable infrastructure — prevents recursion bottlenecks in the Clone-Discard pattern."),
        ],
        notes="One of our biggest innovations was the 'Fair Replay' module. Because HEFT does not see traffic, it looks very fast on paper. We re-run all algorithms' placement plans under the same physical NoC rules to make the comparison fair.",
        fig=FIGURE_MAP[4], fig_desc=FIGURE_DESC[4], right_ratio=0.40,
    ),
    # ── Slide 5 ──────────────────────────────────────────────────────────
    dict(
        slide_num=5,
        title="Deep Dive I — The CA-D Scheduling Pipeline",
        bullets=[
            (0, "**Stage 1: Upward Priority Ranking** — Critical path priority assignment for all DAG tasks."),
            (0, "**Stage 2: Candidate Processor Exploration** — Evaluate all 4×4 grid tiles for each ready task."),
            (0, "**Stage 3: Isolated State Branching** — Clone-Discard pattern prevents data pollution across candidate evaluations."),
            (0, "**Stage 4: Best Candidate Promotion** — Commits the placement that minimises Earliest Finish Time (EFT)."),
            (0, "**Stage 5: Post-Processing Clean-up** — Conservative pruning removes redundant task copies with no live consumers."),
        ],
        notes="Our CA-D algorithm works like a 5-stage pipeline. The 'Clone-Discard' software pattern we invented prevents uncertain trials from permanently polluting the scheduling table.",
        fig=FIGURE_MAP[5], fig_desc=FIGURE_DESC[5], right_ratio=0.42,
    ),
    # ── Slide 6 ──────────────────────────────────────────────────────────
    dict(
        slide_num=6,
        title="Deep Dive II — Recursive Ancestor Duplication Mechanics",
        bullets=[
            (0, "**The Core Evaluation Metric (ΔEFT):**"),
            (1, "*Baseline (EFT_no_dup):* Leave predecessor remote — suffer NoC contention delay."),
            (1, "*Duplicated (EFT_dup):* Replicate predecessor locally on the target core."),
            (1, "*Commit Rule:* ΔEFT = EFT_no_dup − EFT_dup > ε"),
            (0, "**Recursive Traversal:** If local duplication helps, climb the ancestor tree and evaluate grandparents recursively."),
            (0, "**Cycle Guard Constraint:** A visiting set strictly blocks infinite loops during traversal."),
        ],
        notes="There is no blind copying here. If duplicating a parent is mathematically advantageous, we perform a backward recursive search looking at its parents as well.",
        fig=FIGURE_MAP[6], fig_desc=FIGURE_DESC[6], right_ratio=0.38,
    ),
    # ── Slide 7 ──────────────────────────────────────────────────────────
    dict(
        slide_num=7,
        title="Deep Dive III — Best-Instance Selection & Pruning",
        bullets=[
            (0, "**Data Source Selection Rule (priority order):**"),
            (1, "1. Earliest Arrival Time (under physical link contention)"),
            (1, "2. Local Advantage — prefer local copy as tie-breaker"),
            (1, "3. Determinism — smaller processor ID for stable ordering"),
            (0, "**Conservative Redundant Pruning — instance deleted if ALL hold:**"),
            (1, "Not a primary (first-scheduled) node instance"),
            (1, "At least one valid instance of the task remains on the chip"),
            (1, "No local successor or outgoing NoC communication depends on it"),
        ],
        notes="Once scheduling is complete, the Conservative Pruning stage kicks in and safely removes copies that have no remaining demand and contribute nothing to performance.",
        fig=None,
    ),
    # ── Slide 8 ──────────────────────────────────────────────────────────
    dict(
        slide_num=8,
        title="Experimental Setup & Target Workloads",
        bullets=[
            (0, "**960 Comprehensive Experimental Runs** across all scheduler × DAG × CCR combinations."),
            (0, "**4 Target Structural DAG Families:** Chain, Fork, Out-tree, Fork-Join."),
            (0, "**4 CCR Spectrum Points:** 0.1 (compute-bound) → 1.0 → 5.0 → 10.0 (communication-bound)."),
            (0, "**Statistical Distribution:** 20 distinct random seeds per configuration — ensures robust, reproducible averages."),
        ],
        notes="We tested our algorithm with 4 deterministic graph families. Using 4 different CCR ratios and 20 seeds per scenario, we ran a total of 960 large experimental runs.",
        fig=FIGURE_MAP[8], fig_desc=FIGURE_DESC[8], right_ratio=0.50,
    ),
    # ── Slide 9 ── DUAL FIGURE ────────────────────────────────────────────
    dict(
        slide_num=9,
        title="Quantitative Performance — Replayed Speedup & Overhead",
        bullets=[
            (0, "**Core Discovery (Replayed Speedup):** At CCR=10.0, CA-D yields the highest speedup. CD-LS *collapses* — performing worse than HEFT (0.68× slowdown)."),
            (0, "**Illusion Exposed (Replay Overhead):** HEFT and CD-LS show 38–85% optimistic error. CA-D holds a flat 1.0 line — no prediction gap."),
            (0, "**Key Takeaway:** Contention modelling is not optional at high CCR — it is the deciding factor between speedup and degradation."),
        ],
        notes="CD-LS, which duplicates without accounting for network traffic, performed even worse than HEFT at CCR=10. In the Replay Overhead chart, we proved that other algorithms present up to 85% illusory optimism, while CA-D matches real hardware 100%.",
        fig=FIGURE_MAP[9][0], fig_desc=FIGURE_DESC[9],
        dual_fig=True, fig2=FIGURE_MAP[9][1],
        right_ratio=0.50,
    ),
    # ── Slide 10 ─────────────────────────────────────────────────────────
    dict(
        slide_num=10,
        title="Qualitative Performance — Fork-Join Gantt Chart Analysis",
        bullets=[
            (0, "**HEFT Strategy:** Stacks independent tasks onto fewer cores to avoid analytical multi-hop costs. Creates a narrow, tall schedule with poor parallelism."),
            (0, "**CA-D Strategy:** Spatially distributes tasks across the 4×4 mesh. Replicates high-frequency ancestor nodes directly onto target cores. Creates a wide, flat schedule."),
            (0, "**Outcome:** ~20% makespan reduction in dense Fork-Join workloads — achieved by trading a small amount of SRAM for significantly reduced network bandwidth."),
        ],
        notes="Looking at the Gantt chart, we see HEFT piling tasks onto a single processor. CA-D uses the full chip in a balanced way through duplications and compresses the finish time by 20%.",
        fig=FIGURE_MAP[10], fig_desc=FIGURE_DESC[10], right_ratio=0.52,
    ),
    # ── Slide 11 ─────────────────────────────────────────────────────────
    dict(
        slide_num=11,
        title="Discussion & Embedded System Design Trade-offs",
        bullets=[
            (0, "**Computational Cost:** Contention-aware decisions require recursive state cloning — ~14× slower compilation runtime. Chain DAGs trigger maximum recursion depth."),
            (0, "**Memory vs. Bandwidth:** Task duplication reduces network congestion but consumes local instruction memory (SRAM) on each tile."),
            (0, "**Power Balance:** Excessive duplication increases dynamic switching power — making Conservative Pruning mandatory for power-constrained SoC targets."),
            (0, "**Scalability:** The heuristic is fully offline (compile-time) — zero runtime overhead on the physical target."),
        ],
        notes="CA-D makes very smart decisions but the cost is high compilation time and SRAM usage. For embedded systems, we can clearly see how vital the Pruning mechanism we designed is to prevent memory and power waste.",
        fig=FIGURE_MAP[11], fig_desc=FIGURE_DESC[11], right_ratio=0.42,
    ),
    # ── Slide 12 ─────────────────────────────────────────────────────────
    dict(
        slide_num=12,
        title="Conclusions and Future Work",
        bullets=[
            (0, "**Status:** Physical 2D-Mesh NoC models, XY routing, recursive CA-D heuristic, and verification suites are fully implemented and experimentally validated."),
            (0, "**Key Result:** CA-D consistently outperforms all baselines at CCR ≥ 5.0 while maintaining zero replay overhead — proving contention-aware scheduling closes the simulation-to-hardware gap."),
            (0, "**Future Work:**"),
            (1, "Support for heterogeneous computing tiles (GPU / DSP cores)"),
            (1, "Active physical power estimation models"),
            (1, "Online adaptive scheduling variants"),
        ],
        notes="In conclusion, we completed all planned software and testing processes. CA-D achieves its design goals: better makespan at high communication intensity, with provably accurate predictions.",
        fig=None,
    ),
]


# ---------------------------------------------------------------------------
# Main build
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    coverage_report = []

    # Cover
    build_cover_slide(prs)

    for sdata in SLIDES:
        snum  = sdata["slide_num"]
        title = sdata["title"]
        notes = sdata.get("notes", "")

        if sdata.get("type") == "table":
            build_table_slide(
                prs, snum, title,
                intro_bullets=sdata["intro_bullets"],
                table_data=sdata["table_data"],
                notes=notes,
            )
        else:
            fig      = sdata.get("fig")
            fig_desc = sdata.get("fig_desc", "")
            dual     = sdata.get("dual_fig", False)
            fig2     = sdata.get("fig2")
            ratio    = sdata.get("right_ratio", 0.42)

            found = build_content_slide(
                prs, snum, title,
                bullets=sdata["bullets"],
                notes=notes,
                fig_path=fig,
                fig_desc=fig_desc,
                dual_fig=dual,
                fig_path2=fig2,
                right_ratio=ratio,
            )

            if fig is not None:
                if dual:
                    status = "found" if found else "placeholder inserted"
                else:
                    exists = Path(fig).exists() if fig else False
                    status = "found" if exists else "placeholder inserted"
                coverage_report.append((snum, fig_desc, status))

    # Save
    prs.save(str(OUTPUT_PATH))
    print(f"\nSaved: {OUTPUT_PATH}")

    # Sanity check
    prs2 = Presentation(str(OUTPUT_PATH))
    print(f"Sanity check — Slide count: {len(prs2.slides)} (expected 13)")
    pic_count = sum(
        1 for slide in prs2.slides
        for shape in slide.shapes
        if shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
    )
    print(f"Sanity check — Picture shapes: {pic_count}")

    # Figure coverage report
    print("\n── Figure Coverage Report ──────────────────────────────────")
    for snum, desc, status in coverage_report:
        icon = "✓" if status == "found" else "✗"
        print(f"  [{icon}] Slide {snum:02d}: {desc} — {status}")
    if not coverage_report:
        print("  (no figure slots in this build)")
    print("────────────────────────────────────────────────────────────")

    print(f"\nOpen with:  open {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
