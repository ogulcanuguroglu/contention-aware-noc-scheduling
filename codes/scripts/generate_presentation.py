"""
Generate CMP720_CA_D_Scheduling_Presentation.pptx
==================================================
Creates a 12-slide IEEE-style academic presentation + 3 backup slides.
Also writes CMP720_CA_D_Presentation_Outline.md.

Run from repository root:
    python scripts/generate_presentation.py
"""

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIGS = ROOT / "results" / "figures" / "phase21_interpretive"
OUT_PPTX = ROOT / "CMP720_CA_D_Scheduling_Presentation.pptx"
OUT_MD   = ROOT / "CMP720_CA_D_Presentation_Outline.md"

# ---------------------------------------------------------------------------
# Color palette (consistent with project figures)
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLUE   = RGBColor(0x21, 0x66, 0xAC)
GREEN  = RGBColor(0x1A, 0x98, 0x50)
RED    = RGBColor(0xD6, 0x60, 0x4D)
AMBER  = RGBColor(0xCC, 0xBB, 0x44)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
MGRAY  = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0xEE, 0xF2, 0xF7)
LBLUE  = RGBColor(0xD0, 0xE4, 0xF7)
LGREEN = RGBColor(0xD4, 0xED, 0xDA)
LRED   = RGBColor(0xF8, 0xD7, 0xD3)

# ---------------------------------------------------------------------------
# Slide geometry (16:9 widescreen)
# ---------------------------------------------------------------------------
SW  = Inches(13.333)
SH  = Inches(7.5)
TH  = Inches(1.10)   # title bar height
CT  = Inches(1.18)   # content top (just below title bar)
M   = Inches(0.38)   # side margin
CW  = SW - 2 * M     # content width = 12.573"
CH  = SH - CT - Inches(0.12)  # content height ~ 6.2"


# ===========================================================================
# Helper functions
# ===========================================================================

def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def title_bar(slide, text: str, section: str = ""):
    """Navy bar across the full top of the slide with white title text."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, TH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left  = Inches(0.32)
    tf.margin_right = Inches(0.32)
    tf.margin_top   = Inches(0.0)
    tf.margin_bottom = Inches(0.0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if section:
        r0 = p.add_run()
        r0.text = section + "  —  "
        r0.font.size = Pt(11)
        r0.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE8)
        r0.font.name = "Calibri"
        r0.font.bold = False
    r = p.add_run()
    r.text = text
    r.font.size = Pt(24 if len(text) > 55 else 26)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

    # Thin accent line below title bar
    line = slide.shapes.add_shape(1, Inches(0), TH - Inches(0.03), SW, Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def slide_number(slide, n: int, total: int = 12):
    txb = slide.shapes.add_textbox(
        SW - Inches(1.1), SH - Inches(0.30), Inches(1.0), Inches(0.25))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{n} / {total}"
    r.font.size = Pt(10)
    r.font.color.rgb = MGRAY
    r.font.name = "Calibri"


def txbox(slide, left, top, width, height, text, size=18,
          bold=False, color=BLACK, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.0)
    tf.margin_top  = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def bullets(slide, left, top, width, height, items, base_size=17):
    """
    items: list of (indent_level, text, bold, color)
      indent_level 0 = main bullet (u2022)
      indent_level 1 = sub-bullet (en-dash)
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.0)
    tf.margin_top  = Inches(0.0)

    for idx, (lvl, text, bld, col) in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5 if lvl == 0 else 2)
        p.space_after  = Pt(0)

        prefix = ("• " if lvl == 0 else "   – ")
        r = p.add_run()
        r.text = prefix + text
        sz = base_size if lvl == 0 else base_size - 2
        r.font.size  = Pt(sz)
        r.font.bold  = bld
        r.font.color.rgb = col if col else BLACK
        r.font.name  = "Calibri"
    return tb


def colored_box(slide, left, top, width, height, fill, text="",
                text_color=WHITE, size=13, bold=True,
                border_color=None, line_width=Pt(0.75)):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if border_color:
        box.line.color.rgb = border_color
        box.line.width = line_width
    else:
        box.line.fill.background()
    if text:
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left  = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top   = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = text_color
        r.font.name  = "Calibri"
    return box


def arrow_down(slide, cx, top, height=Inches(0.28)):
    """Vertical downward arrow as text."""
    w = Inches(0.4)
    txbox(slide, cx - w / 2, top, w, height,
          "▼", size=11, color=NAVY, align=PP_ALIGN.CENTER)


def add_figure(slide, img_name, left, top, width=None, height=None):
    path = FIGS / img_name
    if not path.exists():
        print(f"  WARNING: figure not found: {path.name}", flush=True)
        return None
    kw = {}
    if width:  kw["width"]  = width
    if height: kw["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kw)


def add_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


# ===========================================================================
# Slide 1 — Title
# ===========================================================================

def slide01_title(prs):
    s = blank_slide(prs)

    # Full navy header block
    hdr = slide.shapes if False else None  # unused

    # Dark background strip at top
    top_bar = s.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(2.8))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY
    top_bar.line.fill.background()

    # Blue accent line
    acc = s.shapes.add_shape(1, Inches(0), Inches(2.8), SW, Inches(0.055))
    acc.fill.solid()
    acc.fill.fore_color.rgb = BLUE
    acc.line.fill.background()

    # Main title
    txbox(s, M, Inches(0.35), CW, Inches(2.2),
          "Contention-Aware Recursive Task Duplication\n"
          "for DAG Scheduling on NoC-Based Multicore Systems",
          size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Course label inside top bar
    txbox(s, M, Inches(0.05), CW, Inches(0.35),
          "CMP720 — Embedded System Design  |  Progress Presentation",
          size=12, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8),
          align=PP_ALIGN.LEFT)

    # Info block below accent line
    info = [
        "Baris Buyukyilmaz",
        "Hacettepe University, Department of Computer Engineering",
        "May 2026",
    ]
    for i, line in enumerate(info):
        sz = 20 if i == 0 else 15
        bld = (i == 0)
        col = BLACK if i > 0 else NAVY
        txbox(s, M, Inches(3.1) + Inches(i * 0.55), CW, Inches(0.55),
              line, size=sz, bold=bld, color=col, align=PP_ALIGN.LEFT)

    # Bottom tag
    txbox(s, M, SH - Inches(0.55), CW, Inches(0.45),
          "Progress presentation — not a final submission",
          size=11, italic=True, color=MGRAY, align=PP_ALIGN.LEFT)

    add_notes(s,
        "Welcome. My name is Baris Buyukyilmaz. This is my progress presentation "
        "for CMP720 Embedded System Design. The project is about scheduling tasks "
        "on a multicore system that uses a Network-on-Chip for communication. "
        "I will explain what the problem is, what I implemented, and what results "
        "I have so far. The presentation is about 12 minutes. Let me begin.")
    return s


# ===========================================================================
# Slide 2 — Problem Definition
# ===========================================================================

def slide02_problem(prs):
    s = blank_slide(prs)
    title_bar(s, "Problem Definition",
              section="1. PROBLEM")
    slide_number(s, 2)

    # Left column: bullets
    bul = [
        (0, "Applications can be modelled as DAGs (Directed Acyclic Graphs)", False, BLACK),
        (1, "Nodes = tasks with computation cost", False, MGRAY),
        (1, "Edges = data dependencies with communication volume", False, MGRAY),
        (0, "Multicore processors are connected via a Network-on-Chip (NoC)", False, BLACK),
        (1, "2D mesh topology with XY routing in this project", False, MGRAY),
        (0, "Remote communication occupies physical NoC links", False, BLACK),
        (0, "Multiple transfers sharing a link cannot overlap", False, BLACK),
        (1, "The second transfer must wait → contention delay", False, RED),
        (0, "Classical schedulers estimate communication cost analytically", False, BLACK),
        (1, "They do not model link sharing or reservation", False, MGRAY),
        (0, "Task duplication can reduce remote data transfers", False, BLACK),
        (1, "But it increases total computation work (TIR > 1.0)", False, MGRAY),
    ]
    bullets(s, M, CT + Inches(0.1), Inches(7.5), CH - Inches(0.1), bul, base_size=16)

    # Right column: key question box
    colored_box(s,
                Inches(8.3), CT + Inches(0.2), Inches(4.6), Inches(1.35),
                fill=NAVY, text="Research Question",
                text_color=WHITE, size=14, bold=True)
    txbox(s, Inches(8.3), CT + Inches(1.6), Inches(4.6), Inches(2.4),
          "Can task duplication reduce makespan\nwhen scheduling on a NoC,\n"
          "if contention is accounted for\nduring the scheduling decision itself?",
          size=15, bold=False, color=NAVY, align=PP_ALIGN.LEFT)

    colored_box(s,
                Inches(8.3), CT + Inches(4.2), Inches(4.6), Inches(0.80),
                fill=LGRAY,
                text="Main message: Scheduling must model both\ncomputation and link-level contention.",
                text_color=NAVY, size=13, bold=False,
                border_color=BLUE)

    add_notes(s,
        "Let me explain the problem. Modern applications can be described as Directed Acyclic Graphs. "
        "Each node is a task, and each edge is a data dependency. We schedule these tasks on multicore "
        "processors that communicate through a Network-on-Chip. The NoC is a 2D mesh in this project. "
        "The key problem is link contention. If two tasks send data at the same time and both paths "
        "share a NoC link, one of them must wait. Classical schedulers like HEFT do not model this. "
        "They use a formula for communication cost but do not reserve links. "
        "Task duplication is one strategy to reduce remote communication. But it adds computation cost. "
        "The question is: can we make better duplication decisions if we know about link contention?")
    return s


# ===========================================================================
# Slide 3 — Proposed Methodology Reminder
# ===========================================================================

def slide03_method_reminder(prs):
    s = blank_slide(prs)
    title_bar(s, "Reminder: Proposed Methodology", section="2. PROPOSAL SUMMARY")
    slide_number(s, 3)

    left_bul = [
        (0, "HEFT baseline: list scheduler, upward-rank priority", True, BLUE),
        (1, "Estimates communication, no link reservations", False, MGRAY),
        (0, "CD-LS extension: add direct parent duplication", True, RED),
        (1, "Duplicate the immediate predecessor onto the target processor", False, MGRAY),
        (1, "Still uses analytic communication model (no link awareness)", False, MGRAY),
        (0, "CA-D proposal: contention-aware recursive duplication", True, GREEN),
        (1, "Evaluate duplication using actual link occupancy", False, MGRAY),
        (1, "Recursively consider ancestor tasks, not just direct parents", False, MGRAY),
        (1, "Accept duplication only if it reduces EFT under contention", False, MGRAY),
    ]
    bullets(s, M, CT + Inches(0.1), Inches(7.5), Inches(4.2), left_bul, base_size=16)

    # Parameter info box
    colored_box(s,
                M, CT + Inches(4.55), Inches(5.8), Inches(1.35),
                fill=LGRAY, border_color=BLUE,
                text="Communication model:   duration = α × hop_count + β × volume\n"
                     "Main experiments:  α = 0,  β = 1  →  duration = volume\n"
                     "NoC: 4×4 homogeneous 2D mesh  |  XY routing  |  16 processors",
                text_color=NAVY, size=13, bold=False)

    # Evaluation box
    colored_box(s,
                Inches(6.6), CT + Inches(4.55), Inches(6.35), Inches(1.35),
                fill=LGRAY, border_color=GREEN,
                text="Evaluation: replay_under_contention()\n"
                     "Preserves task placement → recomputes timing with link reservations\n"
                     "Metric: replayed makespan vs HEFT replayed makespan",
                text_color=NAVY, size=13, bold=False)

    add_notes(s,
        "Before I describe what I actually built, let me briefly remind you of the proposal. "
        "The idea was to compare three schedulers. First, HEFT, which is the standard baseline. "
        "It schedules tasks using an estimated communication cost but does not reserve NoC links. "
        "Second, CD-LS, which adds parent duplication on top of HEFT. "
        "Third, our proposed CA-D, which evaluates duplication using actual link occupancy. "
        "The communication model uses two parameters: alpha for hop cost and beta for volume cost. "
        "In our main experiments, alpha is zero, so duration equals volume. "
        "To evaluate fairly, we use a replay mechanism that recomputes timing with real link reservations. "
        "This gives us a fair comparison across all schedulers. Let me now describe what was actually built.")
    return s


# ===========================================================================
# Slide 4 — Implemented System Architecture
# ===========================================================================

def slide04_architecture(prs):
    s = blank_slide(prs)
    title_bar(s, "Implemented System Architecture", section="3. IMPLEMENTATION")
    slide_number(s, 4)

    BW = Inches(3.30)   # box width
    BH = Inches(0.62)   # box height
    GAP_Y = Inches(0.32)
    CX = SW / 2         # center x of pipeline

    # Row 0: two input boxes side by side
    y0 = CT + Inches(0.10)
    colored_box(s, CX - Inches(4.2), y0, Inches(3.7), BH, fill=BLUE,
                text="DAG Generator\n4 families: Chain, Fork, Out-tree, Fork-join",
                size=12)
    colored_box(s, CX + Inches(0.5), y0, Inches(3.7), BH, fill=BLUE,
                text="MeshNoC Model\n4×4 mesh, XY routing, α=0, β=1",
                size=12)

    # Arrows from both input boxes down to scheduler row
    arrow_down(s, CX - Inches(2.35), y0 + BH, GAP_Y)
    arrow_down(s, CX + Inches(2.35), y0 + BH, GAP_Y)

    # Row 1: three scheduler boxes
    y1 = y0 + BH + GAP_Y
    sched_w = Inches(3.65)
    colored_box(s, M,                y1, sched_w, BH, fill=BLUE,
                text="HEFT\nAnalytic comm, no duplication", size=12)
    colored_box(s, M + sched_w + Inches(0.25), y1, sched_w, BH, fill=RED,
                text="CD-LS\nParent dup, analytic comm", size=12)
    colored_box(s, M + 2 * sched_w + Inches(0.50), y1, sched_w, BH, fill=GREEN,
                text="CA-D\nRecursive ancestor dup, contention-aware", size=12)

    arrow_down(s, CX, y1 + BH, GAP_Y)

    # Row 2: ScheduleState
    y2 = y1 + BH + GAP_Y
    colored_box(s, M, y2, CW, BH, fill=NAVY,
                text="ScheduleState  —  processor_intervals + link_intervals", size=13)

    arrow_down(s, CX, y2 + BH, GAP_Y)

    # Row 3: replay
    y3 = y2 + BH + GAP_Y
    colored_box(s, M, y3, CW, BH, fill=AMBER,
                text="replay_under_contention()  —  re-run timing with link reservations",
                text_color=NAVY, size=13)

    arrow_down(s, CX, y3 + BH, GAP_Y)

    # Row 4: metrics + figures
    y4 = y3 + BH + GAP_Y
    mw = CW / 2 - Inches(0.1)
    colored_box(s, M, y4, mw, BH, fill=LGRAY, border_color=BLUE,
                text="Metrics\nSpeedup, Overhead, TIR, RCVR",
                text_color=NAVY, size=12)
    colored_box(s, M + mw + Inches(0.2), y4, mw, BH, fill=LGRAY, border_color=BLUE,
                text="Figures\n10 PNG+PDF per experiment phase",
                text_color=NAVY, size=12)

    add_notes(s,
        "This slide shows the full framework that was implemented. "
        "At the top, we have two inputs: the DAG generator that creates structured test graphs, "
        "and the MeshNoC model that defines the 4 by 4 processor mesh with XY routing. "
        "All three schedulers take both inputs and produce a ScheduleState object. "
        "The ScheduleState stores two things: processor intervals, which are the task reservations, "
        "and link intervals, which are the communication reservations. "
        "After scheduling, we pass the result to the replay evaluator. "
        "This re-runs the timing with full link-level contention modelling. "
        "Finally, we compute metrics and generate figures. "
        "Everything is implemented in Python using modular source files under the src directory.")
    return s


# ===========================================================================
# Slide 5 — Methodology Updates
# ===========================================================================

def slide05_updates(prs):
    s = blank_slide(prs)
    title_bar(s,
              "Methodology Updates: What Was Applied, What Changed, and Why?",
              section="3. IMPLEMENTATION")
    slide_number(s, 5)

    # Three equal columns with small gaps
    gap   = Inches(0.18)
    col_w = (CW - 2 * gap) / 3          # ~4.06"
    hdr_h = Inches(0.44)                 # header box height
    bul_t = CT + hdr_h + Inches(0.08)   # bullets start y
    bul_h = SH - bul_t - Inches(0.72)   # bullets area height
    c1x   = M
    c2x   = M + col_w + gap
    c3x   = M + 2 * (col_w + gap)

    # ---- Column 1: Applied from Proposal (green) ----------------------------
    colored_box(s, c1x, CT, col_w, hdr_h,
                fill=GREEN, border_color=None,
                text="Applied from Proposal", text_color=WHITE, size=14, bold=True)
    bul1 = [
        (0, "Upward-rank priority (HEFT formula)", False, BLACK),
        (1, "Identical ordering across all three schedulers", False, MGRAY),
        (0, "List scheduling: tasks in rank order", False, BLACK),
        (1, "Ready = all parents scheduled", False, MGRAY),
        (0, "Task duplication to reduce comm wait", False, BLACK),
        (1, "Extended from CD-LS to recursive ancestors", False, MGRAY),
        (0, "4x4 mesh NoC, XY routing", False, BLACK),
        (0, "comm_cost = alpha x hop + beta x vol", False, BLACK),
        (1, "alpha=0 isolates bandwidth contention", False, MGRAY),
    ]
    bullets(s, c1x, bul_t, col_w, bul_h, bul1, base_size=13)

    # ---- Column 2: What Changed (amber) ------------------------------------
    colored_box(s, c2x, CT, col_w, hdr_h,
                fill=AMBER, border_color=None,
                text="What Changed", text_color=WHITE, size=14, bold=True)
    bul2 = [
        (0, "Scalar penalty  ->  per-link reservation", True, RED),
        (1, "link_intervals: sorted list per directed link", False, MGRAY),
        (1, "earliest_route_slot(): first jointly-free slot", False, MGRAY),
        (0, "probe_communication_arrival() added", True, RED),
        (1, "Read-only contention check — no link committed", False, MGRAY),
        (0, "Single critical-parent  ->  all predecessors", True, RED),
        (1, "Every direct parent evaluated independently", False, MGRAY),
        (1, "Greedy recursive: commit if delta-EFT > eps", False, MGRAY),
        (0, "Conservative pruning pass added", False, BLACK),
        (1, "4 safety conditions must all hold for removal", False, MGRAY),
    ]
    bullets(s, c2x, bul_t, col_w, bul_h, bul2, base_size=13)

    # ---- Column 3: Why (blue) ----------------------------------------------
    colored_box(s, c3x, CT, col_w, hdr_h,
                fill=BLUE, border_color=None,
                text="Why", text_color=WHITE, size=14, bold=True)
    bul3 = [
        (0, "Per-link reservation: accurate timing", False, BLACK),
        (1, "Scalar estimate misses link-level congestion", False, MGRAY),
        (1, "Each link independently tracks occupancy", False, MGRAY),
        (0, "probe(): side-effect-free evaluation", False, BLACK),
        (1, "Candidate processor loop clones state", False, MGRAY),
        (1, "Rejected clones are discarded cleanly", False, MGRAY),
        (0, "All-predecessor eval: simpler + deterministic", False, BLACK),
        (1, "Critical-path choice ambiguous under contention", False, MGRAY),
        (1, "All-pred captures more duplication benefit", False, MGRAY),
        (0, "Pruning: correctness safety net", False, BLACK),
        (1, "Rarely triggered; benefit from dup itself", False, MGRAY),
    ]
    bullets(s, c3x, bul_t, col_w, bul_h, bul3, base_size=13)

    # ---- Bottom summary bar ------------------------------------------------
    colored_box(s, M, SH - Inches(0.62), CW, Inches(0.48),
                fill=LGRAY, border_color=NAVY,
                text="Key shift: from a scalar contention estimate to per-link interval "
                     "reservation — enabling CA-D to make duplication decisions under "
                     "realistic link occupancy.",
                text_color=NAVY, size=12, bold=False)

    add_notes(s,
        "This slide organises the methodology into three columns: what was kept from "
        "the proposal, what changed during implementation, and why. "
        "The most important items are in the middle column. "
        "First: the original proposal described a scalar contention penalty. "
        "The implementation uses explicit per-link interval lists instead. "
        "Each directed link on the NoC has its own sorted list of occupied intervals. "
        "When a communication is placed, earliest_route_slot finds the first gap "
        "that is free on every link along the XY route simultaneously. "
        "Second: probe_communication_arrival is a read-only version of this check. "
        "CA-D uses it while evaluating candidate processors without committing anything. "
        "Third: the original idea was to follow Sinnen's critical-parent selection. "
        "The implementation instead evaluates all direct predecessors independently. "
        "This is simpler and deterministic, and captures more duplication benefit. "
        "Finally, a conservative pruning pass was added as a safety net, "
        "but it never triggered on our structured DAG families.")
    return s


# ===========================================================================
# Slide 6 — Scheduler Comparison
# ===========================================================================

def slide06_schedulers(prs):
    s = blank_slide(prs)
    title_bar(s, "Scheduler Comparison: HEFT, CD-LS, CA-D", section="3. IMPLEMENTATION")
    slide_number(s, 6)

    # Compact bullet summary at top
    bul = [
        (0, "HEFT: no duplication, analytic communication, no link reservations", False, BLUE),
        (0, "CD-LS: duplicate direct parent only, still analytic, still contention-blind", False, RED),
        (0, "CA-D: recursive ancestor duplication evaluated under actual link occupancy", False, GREEN),
    ]
    bullets(s, M, CT + Inches(0.08), CW, Inches(0.85), bul, base_size=15)

    # Figure below
    add_figure(s, "fig2_scheduler_concept.png",
               M, CT + Inches(1.0), width=CW)

    add_notes(s,
        "This figure illustrates the conceptual difference between the three schedulers. "
        "I am using a simple chain with three tasks: T0, T1, and T2, each on a different processor. "
        "In panel A, HEFT places each task on the best processor based on estimated finish time. "
        "It does not model link sharing, so two remote communications can appear to overlap. "
        "In panel B, CD-LS duplicates the direct parent. So T1 is duplicated near T2. "
        "This eliminates one communication. But T0 to T2 is still remote. "
        "In panel C, CA-D also duplicates T0. Both ancestors are now local to T2. "
        "All communication on P2 is eliminated. "
        "The key difference of CA-D is that it checks whether this duplication actually helps "
        "under the contention model before committing it.")
    return s


# ===========================================================================
# Slide 7 — Why Replay Is Needed
# ===========================================================================

def slide07_replay(prs):
    s = blank_slide(prs)
    title_bar(s, "Why Replay Is Needed: Exposing Hidden Contention",
              section="3. IMPLEMENTATION")
    slide_number(s, 7)

    bul = [
        (0, "Native schedule: timing computed without link reservations", False, BLACK),
        (0, "replay_under_contention(): recomputes timing with full link-level contention", True, NAVY),
        (1, "Preserves task→processor assignment, discards all original timing", False, MGRAY),
        (1, "Serializes communications that share any NoC link", False, MGRAY),
        (0, "Replay overhead ratio  =  replayed makespan / native makespan", False, BLACK),
        (1, "Ratio = 1.0 → model was accurate.  Ratio > 1.0 → model was optimistic", False, MGRAY),
        (0, "Used as the common fair metric across all schedulers", True, BLUE),
    ]
    bullets(s, M, CT + Inches(0.08), CW, Inches(1.75), bul, base_size=16)

    add_figure(s, "fig3_native_vs_replay.png",
               M, CT + Inches(1.90), width=CW)

    add_notes(s,
        "This slide explains why we need the replay step. "
        "When HEFT or CD-LS builds a schedule, they compute timing analytically. "
        "They do not model what happens when two communications share a link. "
        "The replay mechanism takes the task-to-processor assignments from any scheduler "
        "and recomputes all timing under the contention-aware model. "
        "In the figure, you can see two communications that both start at time zero "
        "in the native model. They use the same link. "
        "Under replay, the second one must wait until the first finishes. "
        "So the actual makespan doubles compared to the prediction. "
        "The replay overhead ratio tells us how optimistic the native model was. "
        "If CA-D's ratio is close to 1.0, it means the scheduler correctly predicted the real behavior. "
        "This is the key evaluation metric we use to compare schedulers fairly.")
    return s


# ===========================================================================
# Slide 8 — DAG Families
# ===========================================================================

def slide08_dag_families(prs):
    s = blank_slide(prs)
    title_bar(s, "DAG Families Used for Evaluation", section="4. RESULTS")
    slide_number(s, 8)

    # Four short descriptions at top (2x2 grid)
    desc = [
        ("Chain",    BLUE,  "Linear dependency chain.\nNegative control: no parallelism."),
        ("Fork",     RED,   "Root → 8 leaves.\nParent dup eliminates all remote comm."),
        ("Out-tree", GREEN, "Depth=2, branching=2.\nRecursive dup provides additional benefit."),
        ("Fork-join",AMBER, "4 branches, each length 3.\nMost sensitive to contention effects."),
    ]
    dw = CW / 4 - Inches(0.12)
    for i, (name, col, desc_text) in enumerate(desc):
        x = M + i * (dw + Inches(0.16))
        colored_box(s, x, CT + Inches(0.08), dw, Inches(0.40),
                    fill=col, text=name, size=13, bold=True)
        txbox(s, x, CT + Inches(0.52), dw, Inches(0.75),
              desc_text, size=12, color=MGRAY, wrap=True)

    add_figure(s, "fig1_dag_family_topologies.png",
               M, CT + Inches(1.35), width=CW)

    add_notes(s,
        "We test the schedulers on four structured DAG families. "
        "Each family stresses a different scheduler behavior. "
        "The chain is the simplest case. It has no useful parallelism. "
        "We use it as a negative control: all schedulers should produce the same result. "
        "The fork DAG has one root and eight leaves. "
        "Here, parent duplication is sufficient. "
        "CA-D and CD-LS behave the same because there is only one ancestor level. "
        "The out-tree has two levels of hierarchy. "
        "This is where recursive ancestor duplication helps CA-D beyond CD-LS. "
        "CA-D can duplicate the root task onto the leaf processors, eliminating all ancestor communication. "
        "The fork-join is the most interesting case. "
        "It has four branches that all reconverge at a join task. "
        "Here, contention between branches becomes very significant, "
        "and contention-blind schedulers may produce worse results under replay.")
    return s


# ===========================================================================
# Slide 9 — Gantt Charts
# ===========================================================================

def slide09_gantt(prs):
    s = blank_slide(prs)
    title_bar(s, "Schedule Examples: Out-tree Gantt Charts", section="4. RESULTS")
    slide_number(s, 9)

    bul = [
        (0, "Out-tree (depth=2, branching=2): shows recursive duplication benefit clearly", False, BLACK),
        (0, "Each row = one processor.  Solid bars = primary tasks.  Hatched bars = duplicate tasks.", False, MGRAY),
        (0, "CCR=1.0 (low): schedulers differ in placement.   CCR=10.0 (high): CA-D eliminates remote comm.", False, BLACK),
    ]
    bullets(s, M, CT + Inches(0.06), CW, Inches(0.78), bul, base_size=15)

    add_figure(s, "fig4b_out_tree_gantt.png",
               M + Inches(0.2), CT + Inches(0.90), width=CW - Inches(0.4))

    add_notes(s,
        "This slide shows the Gantt charts for the out-tree DAG. "
        "Each row is a processor. Each colored bar is a task. "
        "Hatched bars are duplicate task instances. Solid bars are primary instances. "
        "There are two rows of panels: CCR=1.0 at the top and CCR=10.0 at the bottom. "
        "Each row has three panels: HEFT, CD-LS, and CA-D. "
        "At CCR=1.0, the differences are moderate. "
        "At CCR=10.0, you can see a clear change. "
        "HEFT places most tasks on one processor because remote communication becomes very expensive. "
        "CD-LS duplicates direct parents, using a few processors. "
        "CA-D duplicates both the parent and grandparent tasks near the leaves. "
        "This allows all leaf tasks to start immediately with local data. "
        "The makespan difference between HEFT and CA-D is visible on the x-axis. "
        "CA-D achieves a significantly shorter makespan at high CCR.")
    return s


# ===========================================================================
# Slide 10 — Results: CCR Sweep
# ===========================================================================

def slide10_results(prs):
    s = blank_slide(prs)
    title_bar(s, "Results: Replayed Speedup and Replay Overhead vs CCR",
              section="4. RESULTS")
    slide_number(s, 10)

    bul = [
        (0, "Replayed speedup vs HEFT  =  HEFT_replayed / scheduler_replayed  (higher is better)", False, BLACK),
        (0, "Replay overhead ratio  =  replayed / native  (1.0 = model was accurate)", False, BLACK),
        (0, "CCR: Communication-to-Computation Ratio  —  higher CCR = communication-dominated", False, MGRAY),
    ]
    bullets(s, M, CT + Inches(0.05), CW, Inches(0.78), bul, base_size=15)

    half_w = CW / 2 - Inches(0.12)
    add_figure(s, "fig5_ccr_sweep_replayed_speedup.png",
               M, CT + Inches(0.88), width=half_w)
    add_figure(s, "fig6_replay_overhead_ratio.png",
               M + half_w + Inches(0.25), CT + Inches(0.88), width=half_w)

    # Caption
    txbox(s, M, SH - Inches(0.42), CW, Inches(0.35),
          "Left: Speedup of CD-LS and CA-D vs HEFT.  Right: Replay overhead ratio. "
          "Chain: all equal.  Fork: CD-LS = CA-D.  Out-tree/Fork-join: CA-D outperforms CD-LS.",
          size=11, color=MGRAY, italic=True)

    add_notes(s,
        "This is the main result slide. "
        "The left figure shows replayed speedup versus HEFT across CCR values for all four DAG families. "
        "On chain, all schedulers are equal. This confirms the negative control. "
        "On fork, CD-LS and CA-D both achieve around 2.8 times speedup at high CCR. "
        "They are equivalent because the fork has only one ancestor level. "
        "On out-tree, CA-D outperforms CD-LS because recursive ancestor duplication adds the root "
        "onto each leaf processor. "
        "On fork-join, CD-LS can drop below 1.0 at CCR=10, meaning it becomes worse than HEFT "
        "when evaluated fairly. CA-D stays above 1.0 consistently. "
        "The right figure shows replay overhead. "
        "CA-D maintains 1.0 overhead on all topologies. "
        "CD-LS shows 1.26 times overhead on fork-join, revealing that its native prediction was too optimistic. "
        "This confirms that contention-aware scheduling is necessary for accurate makespan prediction.")
    return s


# ===========================================================================
# Slide 11 — Duplication Cost and Communication Reduction
# ===========================================================================

def slide11_duplication(prs):
    s = blank_slide(prs)
    title_bar(s, "Duplication Cost and Remote Communication Reduction",
              section="4. RESULTS")
    slide_number(s, 11)

    bul = [
        (0, "Task Instance Ratio (TIR)  =  total instances / original task count", False, BLACK),
        (1, "TIR = 1.0 means no duplication.  Higher TIR = more computation overhead.", False, MGRAY),
        (0, "Remote Comm Volume Ratio (RCVR)  =  remote edge volume / total DAG volume", False, BLACK),
        (1, "Lower RCVR means more communication was made local by duplication.", False, MGRAY),
        (0, "Trade-off: higher TIR should produce lower RCVR and lower replayed makespan", False, NAVY),
    ]
    bullets(s, M, CT + Inches(0.05), CW, Inches(1.1), bul, base_size=15)

    half_w = CW / 2 - Inches(0.12)
    add_figure(s, "fig7_task_instance_ratio.png",
               M, CT + Inches(1.20), width=half_w)
    add_figure(s, "fig8_remote_comm_volume_ratio.png",
               M + half_w + Inches(0.25), CT + Inches(1.20), width=half_w)

    txbox(s, M, SH - Inches(0.42), CW, Inches(0.35),
          "Left: TIR vs CCR per DAG family.  Right: RCVR vs CCR. "
          "CA-D achieves lower RCVR than CD-LS on out-tree and fork-join.",
          size=11, color=MGRAY, italic=True)

    add_notes(s,
        "This slide shows two metrics that explain the duplication behavior. "
        "The left figure shows the Task Instance Ratio, or TIR. "
        "This measures how many total task instances exist compared to the original task count. "
        "On chain, TIR is 1.0 for all schedulers because duplication does not help there. "
        "On fork and out-tree, CA-D and CD-LS show TIR above 1.0, meaning duplicates were created. "
        "On fork-join, CA-D creates more duplicates than CD-LS, but with better results. "
        "The right figure shows the Remote Communication Volume Ratio. "
        "This is a new metric I introduced. It tells us what fraction of total DAG communication "
        "remains as remote transfers in the native schedule. "
        "Lower is better. "
        "CA-D consistently achieves lower RCVR than CD-LS on out-tree and fork-join. "
        "This means CA-D eliminates more remote communication by placing more ancestor copies locally. "
        "The combination of TIR and RCVR explains why CA-D achieves lower replayed makespan "
        "even though it uses more task instances.")
    return s


# ===========================================================================
# Slide 12 — Discussion, Limitations, Future Work
# ===========================================================================

def slide12_conclusion(prs):
    s = blank_slide(prs)
    title_bar(s, "Discussion, Limitations, and Future Work",
              section="5. CONCLUSIONS")
    slide_number(s, 12)

    col_w = CW / 3 - Inches(0.15)

    # Column headers
    for i, (label, col) in enumerate([
        ("Achievements", GREEN),
        ("Limitations", RED),
        ("Future Work", BLUE),
    ]):
        x = M + i * (col_w + Inches(0.22))
        colored_box(s, x, CT + Inches(0.05), col_w, Inches(0.42),
                    fill=col, text=label, size=14, bold=True)

    # Column content
    ach = [
        (0, "Full DAG/NoC scheduling framework", False, BLACK),
        (0, "HEFT, CD-LS, CA-D implemented", False, BLACK),
        (0, "Explicit link-level contention model", False, BLACK),
        (0, "Fair replay evaluator", False, BLACK),
        (0, "4 DAG families × 4 CCR × 20 seeds", False, BLACK),
        (0, "New metric: RCVR", False, BLACK),
        (0, "Complete figure + report pipeline", False, BLACK),
    ]
    lim = [
        (0, "Synthetic DAGs only", False, BLACK),
        (0, "CA-D is greedy, not optimal", False, BLACK),
        (0, "Not full Sinnen critical-parent", False, RED),
        (0, "Pruning did not trigger in tests", False, BLACK),
        (0, "No cycle-accurate NoC simulation", False, BLACK),
        (0, "Single NoC topology tested", False, BLACK),
        (0, "Runtime grows for deep graphs", False, BLACK),
    ]
    fut = [
        (0, "Finalize figure captions + report", False, BLACK),
        (0, "Add interpretation table to report", False, BLACK),
        (0, "Real application DAGs (e.g., FFT)", False, BLACK),
        (0, "Energy-aware duplication cost", False, BLACK),
        (0, "Larger NoC topologies", False, BLACK),
        (0, "Runtime optimization of CA-D", False, BLACK),
        (0, "Comparison with more algorithms", False, BLACK),
    ]

    for i, col_bullets in enumerate([ach, lim, fut]):
        x = M + i * (col_w + Inches(0.22))
        bullets(s, x, CT + Inches(0.55), col_w, CH - Inches(0.6),
                col_bullets, base_size=14)

    # Bottom summary
    colored_box(s, M, SH - Inches(0.62), CW, Inches(0.47),
                fill=LGRAY, border_color=NAVY,
                text="Main finding: CA-D produces accurate and stable replayed makespans. "
                     "Contention-aware scheduling matters most on multi-level, reconverging DAGs at high CCR.",
                text_color=NAVY, size=13, bold=False)

    add_notes(s,
        "Let me conclude. "
        "On the achievements side, I have a working framework that includes all three schedulers, "
        "a contention-aware replay evaluator, and results across four DAG families and four CCR values. "
        "I also introduced a new metric, the Remote Communication Volume Ratio, "
        "which quantifies how much communication remains remote after duplication. "
        "On the limitations side, the most important is that these are synthetic DAGs. "
        "Results may differ on real application graphs. "
        "Also, CA-D is a greedy heuristic, not an optimal algorithm. "
        "It is not a full implementation of the Sinnen critical-parent algorithm. "
        "The conservative pruning step was implemented but did not trigger in any test case. "
        "For future work, the immediate next step is to finalize the report. "
        "Beyond the course, interesting directions include real application DAGs, "
        "energy-aware scheduling, and larger NoC topologies. "
        "The main finding is that contention-aware scheduling produces stable and accurate results, "
        "especially on topologies where multiple branches share NoC paths. "
        "Thank you. I am ready for questions.")
    return s


# ===========================================================================
# Backup Slide B1 — Communication model details
# ===========================================================================

def slide_b1_comm_model(prs):
    s = blank_slide(prs)
    title_bar(s, "BACKUP — Communication Model Details")

    txbox(s, M, CT + Inches(0.1), CW, Inches(0.38),
          "Communication duration formula:", size=16, bold=True, color=NAVY)
    txbox(s, M, CT + Inches(0.55), CW, Inches(0.50),
          "duration  =  α × hop_count(src, dst)  +  β × communication_volume",
          size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    tbl = [
        ("Parameter", "Value (main)", "Role"),
        ("α (alpha)", "0.0", "Per-hop latency coefficient"),
        ("β (beta)",  "1.0", "Bandwidth coefficient (volume penalty)"),
        ("hop_count", "XY Manhattan distance", "Computed from 4x4 mesh coordinates"),
        ("volume",    "Edge attribute in DAG", "Data transferred across the edge"),
    ]
    row_h = Inches(0.50)
    col_ws = [Inches(2.8), Inches(2.8), Inches(6.5)]
    row_cols = [NAVY, LGRAY, LGRAY, LGRAY, LGRAY]
    row_text_cols = [WHITE, BLACK, BLACK, BLACK, BLACK]
    y_start = CT + Inches(1.15)
    for r, (row, bg, tc) in enumerate(zip(tbl, row_cols, row_text_cols)):
        x = M
        for c, (cell, cw) in enumerate(zip(row, col_ws)):
            colored_box(s, x, y_start + r * row_h, cw, row_h,
                        fill=bg, text=cell, text_color=tc,
                        size=13 if r > 0 else 14, bold=(r == 0),
                        border_color=WHITE)
            x += cw

    txbox(s, M, y_start + len(tbl) * row_h + Inches(0.15), CW, Inches(0.4),
          "Local communication (same processor): duration = 0. No CommunicationInstance created.",
          size=13, color=MGRAY, italic=True)

    bul2 = [
        (0, "With α=0: duration = volume only. Hop distance does not affect timing.", False, BLACK),
        (0, "This isolates bandwidth-driven contention from routing distance effects.", False, BLACK),
        (0, "Alpha sensitivity tested separately: α ∈ {0.0, 1.0, 5.0} (Phase 19).", False, MGRAY),
    ]
    bullets(s, M, y_start + len(tbl) * row_h + Inches(0.65), CW, Inches(1.2),
            bul2, base_size=14)

    add_notes(s, "Backup slide. This shows the communication model in detail.")
    return s


# ===========================================================================
# Backup Slide B2 — CA-D pseudo-code
# ===========================================================================

def slide_b2_pseudocode(prs):
    s = blank_slide(prs)
    title_bar(s, "BACKUP — CA-D Scheduling Logic (Simplified)")

    pseudo = [
        "for each task T in upward-rank order:",
        "   for each candidate processor P:",
        "      state_copy = clone current ScheduleState",
        "      for each direct predecessor pred of T (ascending task_id):",
        "         if pred not already on P:",
        "            EFT_no_dup = contention_EFT(T, P, state_copy)",
        "            _place_recursive_duplicate(pred, P, state_copy)",
        "            EFT_dup    = contention_EFT(T, P, state_copy)",
        "            if EFT_no_dup - EFT_dup <= eps:",
        "               rollback(state_copy)   # no benefit -> reject dup",
        "      candidate_EFT[P] = contention_EFT(T, P, state_copy)",
        "   best_P = argmin(candidate_EFT)",
        "   commit schedule of T on best_P",
        "",
        "_place_recursive_duplicate(pred, P, state):",
        "   if pred already on P: return",
        "   for each grandparent gp of pred:",
        "      EFT_no_gp = EFT of pred on P without gp",
        "      place gp tentatively on P",
        "      EFT_gp    = EFT of pred on P with gp",
        "      if EFT_no_gp - EFT_gp <= eps: rollback gp",
        "      else: commit gp (recursively)",
    ]

    txbox(s, M, CT + Inches(0.05), CW * 0.55, Inches(0.3),
          "CA-D scheduling loop:", size=14, bold=True, color=NAVY)

    code_box = s.shapes.add_textbox(M, CT + Inches(0.40), CW, SH - CT - Inches(0.55))
    tf = code_box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.1)
    tf.margin_top  = Inches(0.05)
    for i, line in enumerate(pseudo):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.name = "Courier New"
        if line.strip().startswith("_place") or line.strip().startswith("for each task"):
            r.font.bold = True
            r.font.color.rgb = NAVY
        elif "EFT_no_dup" in line or "EFT_dup" in line or "rollback" in line:
            r.font.color.rgb = RED
        elif "EFT_gp" in line or "gp" in line:
            r.font.color.rgb = GREEN
        else:
            r.font.color.rgb = BLACK

    add_notes(s, "Backup slide. Simplified pseudo-code for the CA-D algorithm.")
    return s


# ===========================================================================
# Backup Slide B3 — Metric definitions
# ===========================================================================

def slide_b3_metrics(prs):
    s = blank_slide(prs)
    title_bar(s, "BACKUP — Metric Definitions")

    metrics = [
        ("Replayed speedup vs HEFT",
         "HEFT_replayed_makespan  /  scheduler_replayed_makespan",
         "Primary comparison metric. Fair because both use contention-aware replay.",
         GREEN),
        ("Replay overhead ratio",
         "replayed_makespan  /  native_makespan",
         "Measures native-model optimism. 1.0 = accurate. >1.0 = scheduler was optimistic.",
         AMBER),
        ("Task Instance Ratio (TIR)",
         "total_task_instances  /  original_task_count",
         "Measures duplication overhead. 1.0 = no duplication. Higher = more copies.",
         BLUE),
        ("Remote Comm Volume Ratio (RCVR)",
         "remote_edge_volume  /  total_DAG_edge_volume",
         "Fraction of communication that remains remote. "
         "Edge u→v is remote if no instance of u is co-located with v's primary processor.",
         RED),
        ("CCR  (Communication-to-Computation Ratio)",
         "total_comm_volume  /  total_computation_cost",
         "Controls experiment regime. Low CCR = compute-dominated. High CCR = comm-dominated.",
         NAVY),
    ]

    y = CT + Inches(0.08)
    for name, formula, desc, col in metrics:
        colored_box(s, M, y, Inches(4.0), Inches(0.35),
                    fill=col, text=name, size=12, bold=True)
        txbox(s, M + Inches(4.15), y, CW - Inches(4.15), Inches(0.35),
              formula, size=12, bold=True, color=col)
        txbox(s, M, y + Inches(0.38), CW, Inches(0.32),
              desc, size=11, color=MGRAY, italic=True)
        y += Inches(0.80)

    add_notes(s, "Backup slide. Formal definitions of all metrics used in the evaluation.")
    return s


# ===========================================================================
# Main
# ===========================================================================

def build_presentation():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Building slides ...", flush=True)
    slide01_title(prs)
    print("  01 Title", flush=True)
    slide02_problem(prs)
    print("  02 Problem Definition", flush=True)
    slide03_method_reminder(prs)
    print("  03 Methodology Reminder", flush=True)
    slide04_architecture(prs)
    print("  04 System Architecture", flush=True)
    slide05_updates(prs)
    print("  05 Methodology Updates", flush=True)
    slide06_schedulers(prs)
    print("  06 Scheduler Comparison", flush=True)
    slide07_replay(prs)
    print("  07 Why Replay", flush=True)
    slide08_dag_families(prs)
    print("  08 DAG Families", flush=True)
    slide09_gantt(prs)
    print("  09 Gantt Charts", flush=True)
    slide10_results(prs)
    print("  10 Results - CCR Sweep", flush=True)
    slide11_duplication(prs)
    print("  11 Duplication Cost", flush=True)
    slide12_conclusion(prs)
    print("  12 Discussion + Future Work", flush=True)
    slide_b1_comm_model(prs)
    print("  B1 Backup - Comm Model", flush=True)
    slide_b2_pseudocode(prs)
    print("  B2 Backup - Pseudo-code", flush=True)
    slide_b3_metrics(prs)
    print("  B3 Backup - Metric Definitions", flush=True)

    prs.save(str(OUT_PPTX))
    print(f"\nSaved -> {OUT_PPTX.name}", flush=True)
    return prs


def write_outline():
    text = """# CMP720 CA-D Scheduling — Presentation Outline

## Slide Count: 12 main + 3 backup = 15 total
## Estimated time: 11–13 minutes (main) + 5 min Q&A

---

## Slide 01 — Title
**Section:** —
**Message:** Project overview: contention-aware recursive duplication on NoC.
**Notes:** ~30 seconds. Welcome and introduce yourself.

---

## Slide 02 — Problem Definition
**Section:** 1. Problem
**Message:** Scheduling on NoC must model both computation and link-level contention.
**Key points:**
- DAGs: tasks + data dependencies
- NoC: physical links shared by multiple communications
- Link contention = serialization delay
- Classical schedulers ignore link sharing
- Task duplication can help but adds overhead
**Notes:** ~75 seconds.

---

## Slide 03 — Proposed Methodology Reminder
**Section:** 2. Proposal Summary
**Message:** Combine task duplication with contention-aware NoC scheduling.
**Key points:**
- HEFT baseline: analytic comm, no link reservation
- CD-LS: parent duplication, still analytic
- CA-D proposed: contention-aware recursive ancestor duplication
- Evaluation via replay mechanism
**Notes:** ~75 seconds.

---

## Slide 04 — Implemented System Architecture
**Section:** 3. Implementation
**Message:** A complete simulation framework was implemented.
**Key points:**
- DAG Generator (4 families)
- MeshNoC model (4x4, XY routing)
- Three scheduler modules
- ScheduleState (processor + link intervals)
- replay_under_contention()
- Metrics and figure pipeline
**Notes:** ~60 seconds.

---

## Slide 05 — Methodology Updates
**Section:** 3. Implementation
**Message:** Contention modelled as explicit link-interval reservation, not a scalar penalty.
**Key points:**
- Scalar penalty → explicit link_intervals dict
- earliest_route_slot() finds joint-free slot on all route links
- probe_communication_arrival() for read-only CA-D evaluation
- CA-D is NOT Sinnen critical-parent selection
- Pruning implemented but did not trigger significantly
**Notes:** ~90 seconds.

---

## Slide 06 — Scheduler Comparison: HEFT, CD-LS, CA-D
**Section:** 3. Implementation
**Message:** CA-D combines recursive duplication with contention-aware route reservation.
**Figure:** fig2_scheduler_concept.png
**Notes:** ~75 seconds.

---

## Slide 07 — Why Replay Is Needed
**Section:** 3. Implementation
**Message:** Replay exposes hidden contention in contention-blind schedules.
**Figure:** fig3_native_vs_replay.png
**Key points:**
- Native model may allow simultaneous link use
- Replay serializes transfers on shared links
- Overhead ratio = replayed / native
**Notes:** ~60 seconds.

---

## Slide 08 — DAG Families Used for Evaluation
**Section:** 4. Results
**Message:** Different DAG structures stress different scheduler behaviors.
**Figure:** fig1_dag_family_topologies.png
**Key points:**
- Chain: negative control (no parallelism)
- Fork: parent dup sufficient
- Out-tree: recursive dup beneficial
- Fork-join: most sensitive to contention
**Notes:** ~60 seconds.

---

## Slide 09 — Schedule Examples: Out-tree Gantt Charts
**Section:** 4. Results
**Message:** CA-D changes placement by putting ancestor copies near consuming tasks.
**Figure:** fig4b_out_tree_gantt.png
**Notes:** ~75 seconds.

---

## Slide 10 — Results: CCR Sweep and Replay Overhead
**Section:** 4. Results
**Message:** CA-D is most useful when recursive duplication and contention-awareness both matter.
**Figures:** fig5_ccr_sweep_replayed_speedup.png + fig6_replay_overhead_ratio.png
**Key observations:**
- Chain: all equal
- Fork: CD-LS = CA-D
- Out-tree: CA-D > CD-LS
- Fork-join: CD-LS can drop below 1.0; CA-D stable
**Notes:** ~90 seconds.

---

## Slide 11 — Duplication Cost and Remote Communication Reduction
**Section:** 4. Results
**Message:** CA-D trades extra task instances for reduced remote communication.
**Figures:** fig7_task_instance_ratio.png + fig8_remote_comm_volume_ratio.png
**Notes:** ~75 seconds.

---

## Slide 12 — Discussion, Limitations, and Future Work
**Section:** 5. Conclusions
**Message:** Functional results, honest limitations, clear next steps.
**3 columns:**
- Achievements: full framework, fair replay, new RCVR metric, 4 DAG x 4 CCR x 20 seeds
- Limitations: synthetic DAGs, greedy heuristic, not full Sinnen, pruning not triggered
- Future Work: real DAGs, energy cost, larger NoC, more algorithms
**Notes:** ~90 seconds.

---

## BACKUP Slide B1 — Communication Model Details
Communication duration formula, parameter table, alpha sensitivity note.

## BACKUP Slide B2 — CA-D Pseudo-code
Simplified pseudo-code for CA-D scheduling loop and recursive ancestor placement.

## BACKUP Slide B3 — Metric Definitions
TIR, RCVR, speedup, overhead, CCR — formal definitions.

---

## Figures Used

| Slide | Figure file |
|-------|------------|
| 6 | fig2_scheduler_concept.png |
| 7 | fig3_native_vs_replay.png |
| 8 | fig1_dag_family_topologies.png |
| 9 | fig4b_out_tree_gantt.png |
| 10 | fig5_ccr_sweep_replayed_speedup.png |
| 10 | fig6_replay_overhead_ratio.png |
| 11 | fig7_task_instance_ratio.png |
| 11 | fig8_remote_comm_volume_ratio.png |

All from: results/figures/phase21_interpretive/

## Figures omitted (available for backup discussion)
- fig4a_fork_gantt.png — fork Gantt (less informative than out-tree for this presentation)
- fig4c_fork_join_gantt.png — fork-join Gantt (discussed verbally in slide 10/12)

## Estimated presentation time
- Slides 1–12: ~11–13 minutes
- Q&A: up to 5 minutes
- Backup slides: on demand during Q&A
"""
    OUT_MD.write_text(text, encoding="utf-8")
    print(f"Saved -> {OUT_MD.name}", flush=True)


if __name__ == "__main__":
    build_presentation()
    write_outline()
    print("\nDone.", flush=True)
    print(f"  PPTX -> {OUT_PPTX}", flush=True)
    print(f"  MD   -> {OUT_MD}", flush=True)
